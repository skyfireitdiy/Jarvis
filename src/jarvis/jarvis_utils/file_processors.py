from pathlib import Path
import fitz  # PyMuPDF for PDF files
from docx import Document as DocxDocument  # python-docx for DOCX files
from pptx import Presentation
import pandas as pd
import unicodedata

class FileProcessor:
    """Base class for file processor"""
    @staticmethod
    def can_handle(file_path: str) -> bool:
        """Determine if the file can be processed"""
        raise NotImplementedError

    @staticmethod
    def extract_text(file_path: str) -> str:
        """Extract file text content"""
        raise NotImplementedError

class TextFileProcessor(FileProcessor):
    """Text file processor"""
    ENCODINGS = ['utf-8', 'gbk', 'gb2312', 'latin1']
    SAMPLE_SIZE = 8192  # Read the first 8KB to detect encoding

    @staticmethod
    def can_handle(file_path: str) -> bool:
        """Determine if the file is a text file by trying to decode it"""
        try:
            # Read the first part of the file to detect encoding
            with open(file_path, 'rb') as f:
                sample = f.read(TextFileProcessor.SAMPLE_SIZE)

            # Check if it contains null bytes (usually represents a binary file)
            if b'\x00' in sample:
                return False

            # Check if it contains too many non-printable characters (usually represents a binary file)
            non_printable = sum(1 for byte in sample if byte < 32 and byte not in (9, 10, 13))  # tab, newline, carriage return
            if non_printable / len(sample) > 0.3:  # If non-printable characters exceed 30%, it is considered a binary file
                return False

            # Try to decode with different encodings
            for encoding in TextFileProcessor.ENCODINGS:
                try:
                    sample.decode(encoding)
                    return True
                except UnicodeDecodeError:
                    continue

            return False

        except Exception:
            return False

    @staticmethod
    def extract_text(file_path: str) -> str:
        """Extract text content, using the detected correct encoding"""
        detected_encoding = None
        try:
            # First try to detect encoding
            with open(file_path, 'rb') as f:
                raw_data = f.read()

            # Try different encodings
            for encoding in TextFileProcessor.ENCODINGS:
                try:
                    raw_data.decode(encoding)
                    detected_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue

            if not detected_encoding:
                raise UnicodeDecodeError(f"Failed to decode file with supported encodings: {file_path}") # type: ignore

            # Use the detected encoding to read the file
            with open(file_path, 'r', encoding=detected_encoding, errors='ignore') as f:
                content = f.read()

            # Normalize Unicode characters
            content = unicodedata.normalize('NFKC', content)

            return content

        except Exception as e:
            raise Exception(f"Failed to read file: {str(e)}")

class PDFProcessor(FileProcessor):
    """PDF file processor"""
    @staticmethod
    def can_handle(file_path: str) -> bool:
        return Path(file_path).suffix.lower() == '.pdf'

    @staticmethod
    def extract_text(file_path: str) -> str:
        """提取PDF文件中的所有文本内容，包括页码、图片描述等"""
        try:
            text_parts = []
            with fitz.open(file_path) as doc:  # type: ignore
                # 添加文档信息
                info = doc.metadata
                if info:
                    meta_text = []
                    if info.get("title"):
                        meta_text.append(f"标题: {info['title']}")
                    if info.get("author"):
                        meta_text.append(f"作者: {info['author']}")
                    if info.get("subject"):
                        meta_text.append(f"主题: {info['subject']}")
                    if info.get("keywords"):
                        meta_text.append(f"关键词: {info['keywords']}")
                    
                    if meta_text:
                        text_parts.append("=== 文档信息 ===")
                        text_parts.append("\n".join(meta_text))
                
                # 提取目录结构（如果有）
                toc = doc.get_toc()  # type: ignore
                if toc:
                    text_parts.append("\n=== 目录结构 ===")
                    for level, title, page in toc:
                        indent = "  " * (level - 1)
                        text_parts.append(f"{indent}- {title} (第{page}页)")
                
                # 处理各页内容
                text_parts.append("\n=== 页面内容 ===")
                for page_index in range(len(doc)):  # 使用范围遍历而不是直接枚举文档对象
                    # 添加页码标记
                    text_parts.append(f"\n--- 第{page_index+1}页 ---")
                    
                    # 获取页面
                    page = doc[page_index]
                    
                    # 提取页面文本（包括结构信息）
                    try:
                        # 尝试使用结构化提取（保留段落和块结构）
                        text = page.get_text("text")  # type: ignore
                        text = text.strip()
                        if text:
                            text_parts.append(text)
                    except Exception:
                        # 如果结构化提取失败，回退到简单文本提取
                        text = page.get_text()  # type: ignore
                        if text.strip():
                            text_parts.append(text.strip())
                    
                    # 提取图像信息（如果需要）
                    # 注意：这可能会增加处理时间，可根据需要启用
                    """
                    image_list = page.get_images()
                    if image_list:
                        text_parts.append(f"本页包含 {len(image_list)} 个图像")
                    """
            
            # 合并所有文本
            return "\n".join(text_parts)
        
        except Exception as e:
            # 处理可能的异常
            return f"PDF处理错误: {str(e)}"

class DocxProcessor(FileProcessor):
    """DOCX file processor"""
    @staticmethod
    def can_handle(file_path: str) -> bool:
        return Path(file_path).suffix.lower() == '.docx'

    @staticmethod
    def extract_text(file_path: str) -> str:
        """提取 DOCX 文件中的所有文本内容，包括段落、表格、页眉页脚等"""
        doc = DocxDocument(file_path)
        full_text = []
        
        # 提取段落文本
        for para in doc.paragraphs:
            if para.text.strip():  # 跳过空段落
                full_text.append(para.text)
        
        # 提取表格文本
        for table in doc.tables:
            for row in table.rows:
                row_texts = []
                for cell in row.cells:
                    # 每个单元格可能包含多个段落
                    cell_text = "\n".join([p.text for p in cell.paragraphs if p.text.strip()])
                    if cell_text:
                        row_texts.append(cell_text)
                if row_texts:
                    full_text.append(" | ".join(row_texts))
        
        # 提取页眉页脚（如果有节）
        try:
            for section in doc.sections:
                # 提取页眉
                if section.header:
                    header_text = "\n".join([p.text for p in section.header.paragraphs if p.text.strip()])
                    if header_text:
                        full_text.append(f"页眉: {header_text}")
                
                # 提取页脚
                if section.footer:
                    footer_text = "\n".join([p.text for p in section.footer.paragraphs if p.text.strip()])
                    if footer_text:
                        full_text.append(f"页脚: {footer_text}")
        except:
            # 如果提取页眉页脚失败，忽略错误继续
            pass
        
        # 合并所有文本
        return "\n\n".join(full_text)

class PPTProcessor(FileProcessor):
    """PPT file processor"""
    @staticmethod
    def can_handle(file_path: str) -> bool:
        return Path(file_path).suffix.lower() in ['.ppt', '.pptx']

    @staticmethod
    def extract_text(file_path: str) -> str:
        """提取PPT文件中的所有文本内容，包括标题、文本框、备注等"""
        prs = Presentation(file_path)
        all_text = []
        
        # 遍历所有幻灯片
        for slide_index, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # 添加幻灯片编号
            slide_text.append(f"=== 幻灯片 {slide_index} ===")
            
            # 提取幻灯片中所有形状的文本
            for shape in slide.shapes:
                # 提取带有文本的形状
                try:
                    if hasattr(shape, "text_frame") and shape.text_frame:  # type: ignore
                        for paragraph in shape.text_frame.paragraphs:  # type: ignore
                            text = paragraph.text.strip()
                            if text:
                                slide_text.append(text)
                except AttributeError:
                    pass
                
                # 提取表格内容
                try:
                    if hasattr(shape, "table") and shape.table:  # type: ignore
                        for row in shape.table.rows:  # type: ignore
                            row_texts = []
                            for cell in row.cells:
                                if hasattr(cell, "text_frame") and cell.text_frame:
                                    cell_paragraphs = cell.text_frame.paragraphs  # type: ignore
                                    cell_text = " ".join([p.text.strip() for p in cell_paragraphs if p.text.strip()])
                                    if cell_text:
                                        row_texts.append(cell_text)
                            if row_texts:
                                slide_text.append(" | ".join(row_texts))
                except AttributeError:
                    pass
            
            # 提取幻灯片备注
            try:
                if hasattr(slide, "has_notes_slide") and slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide and hasattr(notes_slide, "notes_text_frame") and notes_slide.notes_text_frame:
                        notes_text = notes_slide.notes_text_frame.text.strip()  # type: ignore
                        if notes_text:
                            slide_text.append(f"备注: {notes_text}")
            except AttributeError:
                pass
            
            # 合并当前幻灯片的所有文本
            if len(slide_text) > 1:  # 如果除了幻灯片编号外还有其他内容
                all_text.append("\n".join(slide_text))
        
        # 返回所有幻灯片的文本内容
        return "\n\n".join(all_text)

class ExcelProcessor(FileProcessor):
    """Excel file processor"""
    @staticmethod
    def can_handle(file_path: str) -> bool:
        return Path(file_path).suffix.lower() in ['.xls', '.xlsx']

    @staticmethod
    def extract_text(file_path: str) -> str:
        """提取 Excel 文件中的所有文本内容，包括多个工作表及格式化内容"""
        try:
            # 读取所有工作表
            excel_file = pd.ExcelFile(file_path)
            sheets_text = []
            
            # 处理每个工作表
            for sheet_name in excel_file.sheet_names:
                # 读取当前工作表
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # 如果是空表格，跳过
                if df.empty:
                    continue
                
                # 添加工作表标题
                sheet_text = [f"=== 工作表: {sheet_name} ==="]
                
                # 填充空单元格，避免NaN显示
                df = df.fillna("")
                
                # 提取表格头信息
                if not df.columns.empty:
                    headers = [str(col) for col in df.columns]
                    sheet_text.append("列标题: " + " | ".join(headers))
                
                # 尝试提取表格中可能的关键信息
                # 1. 表格内容概述
                row_count, col_count = df.shape
                sheet_text.append(f"表格大小: {row_count}行 x {col_count}列")
                
                # 2. 表格数据，使用更友好的格式
                try:
                    # 转换数据框为字符串表示
                    # 设置最大行数和列数，避免过大的表格
                    max_rows = min(500, row_count)  # 最多显示500行
                    if row_count > max_rows:
                        sheet_text.append(f"注意: 表格太大，仅显示前{max_rows}行")
                    
                    # 将DataFrame转换为字符串表格
                    table_str = df.head(max_rows).to_string(index=True, max_rows=max_rows, max_cols=None)
                    sheet_text.append(table_str)
                    
                except Exception as e:
                    sheet_text.append(f"表格数据提取错误: {str(e)}")
                
                # 合并当前工作表的文本
                sheets_text.append("\n".join(sheet_text))
            
            # 如果没有提取到任何内容，返回一个提示信息
            if not sheets_text:
                return "Excel文件为空或无法提取内容"
                
            # 合并所有工作表的文本
            return "\n\n".join(sheets_text)
            
        except Exception as e:
            # 处理可能的异常，返回错误信息
            return f"Excel文件处理错误: {str(e)}"
